# -*- coding: utf-8 -*-
# Copyright Dr. Jahn von Heys

""" Writes ""Hallo Welt"" to Excel, Powerpoint, Access file.

An advanced version with class and methods.
"""

# Offizielle Python Pakete / Module in alphabetischer Reihenfolge
import openpyxl      # Bibliothek zum Lesen und Schreiben von Excel
import pptx          # Bibliothek zum Lesen und Schreiben von Powerpoint
import pyodbc        # Bibliothek zum Lesen und Schreiben von Datenbanken (SQL)

# HCOB-Module in alphabetischer Reihenfolge
import msaccessdb  # Erstellen einer Access DB


HELLO_WORLD_MESSAGE = 'Hallo Welt!'


class MultiPrinter():
    """ This class demonstrates methods, attributes and constants.
    """
    def __init__(self, message=None, xls_file="HalloWelt.xlsx",
                 db_file=r"C:\Service\Development\Libs\python_libs\hcob-packages\python_examples\HalloWelt.accdb",
                 ppt_file="HalloWelt.pptx"):
        self.message = message or HELLO_WORLD_MESSAGE
        self.xls_file = xls_file
        self.db_file = db_file
        self.ppt_file = ppt_file

    def write_to_console(self):
        """ Prints the message attribute on console.
        """
        print(self.message)

    def write_to_excel(self):
        """ Writes message attribute to Excel file.

        An existing file is overwritten without warning!

        """
        print("Write to Excel...")
        wb = openpyxl.Workbook()  # Neues Tabellendokument
        ws = wb['Sheet']          # Tabellenblatt "Sheet" auswählen
        ws['A1'] = self.message.split()[0]
        ws['A2'] = self.message.split()[1]
        ws['A3'] = '=A1 & " " & A2'
        wb.save(self.xls_file)  # Tabellendokument speichern

    def write_to_access(self):
        """ Writes message attribute to Access DB.

        Access DB is created if it doesn't exist.

        """
        print("Write to Access DB...")
        msaccessdb.create(self.db_file)
        conn_str = 'Driver={Microsoft Access Driver (*.mdb, *.accdb)};DBQ=' + \
            self.db_file + ';'
        conn = pyodbc.connect(conn_str)
        cursor = conn.cursor()
        sql = \
            """
            CREATE TABLE tbl_HalloWelt
            ( Spalte1 Text,
              Spalte2 Text
            )
            """

        cursor.execute(sql)
        conn.commit()

        sql = \
            """
            INSERT INTO tbl_HalloWelt
            VALUES
            ('{0}', '{1}');
            """.format(self.message.split()[0], self.message.split()[1])

        cursor.execute(sql)
        conn.commit()

    def write_to_powerpoint(self):
        """ Writes message attribute to powerpoint file.
        """
        print("Write to Powerpoint...")
        prs = pptx.Presentation()                           # Neue Präsentation
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Neue Slide
        slide.shapes.title.text = self.message              # Slide-Titel
        prs.save(self.ppt_file)                             # Präsi speichern


if __name__ == '__main__':
    # Auf geht's: Instanz mit Standardwerten
    multi_printer_default = MultiPrinter()
    multi_printer_default.write_to_console()
    multi_printer_default.write_to_excel()
    multi_printer_default.write_to_access()
    multi_printer_default.write_to_powerpoint()

    # Nun individuell
    multi_printer_custom = MultiPrinter(message='Hello world!',
                                        xls_file='CustomHello.xlsx',
                                        db_file=r'C:\Service\Development\Libs\python_libs\hcob-packages\python_examples\CustomHello.accdb',
                                        ppt_file='CustomHello.pptx')
    multi_printer_custom.write_to_console()
    multi_printer_custom.write_to_excel()
    multi_printer_custom.write_to_access()
    multi_printer_custom.write_to_powerpoint()

    multi_printer_default.new_message = 'New message'
    print(multi_printer_default.new_message)

    print('Finished.')
