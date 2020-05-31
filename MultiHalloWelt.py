# -*- coding: utf-8 -*-
"""
Created on Tue Aug 13 14:46:18 2019

@author: U0045119
"""
# Bibliothek zum Lesen und Schreiben von Datenbanken (SQL)
import pyodbc
import msaccessdb

# Bibliothek zum Lesen und Schreiben von Excel
import openpyxl
# Bibliothek zum Lesen und Schreiben von Powerpoint
import pptx

# Ausgabe auf der Console #####################################################
print("Hallo Welt!")

# Ausgabe in Excel ############################################################
print("Write to Excel...")
wb = openpyxl.Workbook()  # Neues Tabellendokument
ws = wb['Sheet']          # Tabellenblatt "Sheet" auswählen
ws['A1'] = "Hallo "
ws['A2'] = "Welt!"
ws['A3'] = "=A1&A2"
wb.save("HalloWelt.xlsx") # Tabellendokument speichern

# Ausgabe in Powerpoint #######################################################
print("Write to Powerpoint...")
prs = pptx.Presentation()                           # Neue Präsentation
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Neue Slide
slide.shapes.title.text = "Hallo Welt!"             # Slide-Titel setzen
prs.save("HalloWelt.pptx")                          # Präsi speichern

# Ausgabe in ACCESS ###########################################################
print("Write to Access DB...")
access_db = 'C:\Service\Development\Libs\python_libs\hcob-packages\python_examples\HalloWelt.accdb'
msaccessdb.create(access_db)
conn = pyodbc.connect(r'Driver={Microsoft Access Driver (*.mdb, *.accdb)};DBQ=' +
                      access_db + ';')
cursor = conn.cursor()
sql = \
"""
CREATE TABLE tbl_HalloWelt
( Spalte1 Text,
  Spalte2 Text
)
"""

cursor.execute(sql)
conn.commit()

sql = \
"""
INSERT INTO tbl_HalloWelt
VALUES
('Hallo', 'Welt');
"""

cursor.execute(sql)
conn.commit()

print('Finished')